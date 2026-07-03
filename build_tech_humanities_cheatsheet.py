from __future__ import annotations

import os
import re
import sys
import textwrap
from collections import Counter, defaultdict
from pathlib import Path

import pdfplumber
from docx import Document
from pptx import Presentation
from reportlab.lib import colors
from reportlab.lib.enums import TA_LEFT
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.pdfbase.pdfmetrics import registerFont
from reportlab.platypus import (
    BaseDocTemplate,
    Frame,
    FrameBreak,
    KeepTogether,
    PageBreak,
    PageTemplate,
    Paragraph,
    Spacer,
    Table,
    TableStyle,
)


SOURCE_DIR = Path(r"D:\科技与人文")
OUT_DIR = Path("output/pdf")
TMP_DIR = Path("tmp/pdfs/tech_humanities")
OUT_PDF = OUT_DIR / "科技与人文_双面A4超密cheatpaper_纯知识点版.pdf"
NOTES_TXT = TMP_DIR / "extracted_source_notes.txt"


FILES = [
    "1-s2.0-S0142961219308294-main.pdf",
    "1-s2.0-S1525001616453323-main.pdf",
    "AI Will Kill Love.docx",
    "Anthropic_ConstitutionalAI_v2(1).pdf",
    "Christopher Reeve Senate Hearing Statement.pdf",
    "Clock：The Ribbon of Time(2)(1).pptx",
    "Controversial Chinese scientist He Jiankui proposes new gene editing research _ CNN.pdf",
    "CRISPR(1).pptx",
    "Group Discussion on Mdedvi.docx",
    "He Jiankui_ Chinese gene-editing scientist jailed for 3 years _ CNN.pdf",
    "How human gene editing is moving on after CRISPR baby scandal _ CNN.pdf",
    "Human Fertilisation and Embryology Act 1990.pdf",
    "LLM_to_AI_Agent_TED_like_deck_v4_visual_minimal_english(1).pptx",
    "Material(1).docx",
    "NFC 科技与人文 .pptx",
    "OPC.pptx",
    "Plastic.pptx",
    "Slides for pre.pptx",
    "social credit system.pptx",
    "Therapeutic Cloning Act 2001.pdf",
    "Therapeutic Cloning Ethics Tribunal Layout.pdf",
    "Therapeutic cloning promises and issues.pdf",
    "案例分析更新.pptx",
    "基于近场通信与射频识别技术的非接触式支付研究_高梦谣.pdf",
    "手机NFC技术的滥用风险与防范_沈臻懿.pdf",
    "余数生命悖论(2).pdf",
    "主权与治理：阿甘本与福柯“生命政治”的模式差异与内在联结_郑雨晨(1).pdf",
    "艰难时世节选.docx",
]


STOPWORDS = set(
    """
    the and for that with from this are was were have has had can will would could should about into
    technology technologies human humans humanities science scientific ethical ethics social society
    china chinese paper research study article presentation group discussion student students slide
    以及 一个 一种 通过 可以 进行 研究 技术 科技 人文 问题 影响 社会 中国 我们 这个 这些 他们 由于 因此
    """.split()
)


def clean_text(text: str) -> str:
    text = text.replace("\x00", " ")
    text = re.sub(r"[ \t\r\f\v]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def extract_pdf(path: Path) -> str:
    chunks = []
    try:
        with pdfplumber.open(str(path)) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text(x_tolerance=1, y_tolerance=3) or ""
                if text.strip():
                    chunks.append(f"[p.{i + 1}] {text}")
    except Exception as exc:
        chunks.append(f"[EXTRACT_ERROR] {exc}")
    return clean_text("\n".join(chunks))


def extract_docx(path: Path) -> str:
    chunks = []
    try:
        doc = Document(str(path))
        for p in doc.paragraphs:
            if p.text.strip():
                chunks.append(p.text)
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    chunks.append(row_text)
    except Exception as exc:
        chunks.append(f"[EXTRACT_ERROR] {exc}")
    return clean_text("\n".join(chunks))


def extract_pptx(path: Path) -> str:
    chunks = []
    try:
        prs = Presentation(str(path))
        for idx, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    texts.append(clean_text(shape.text))
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                        if row_text:
                            texts.append(row_text)
            if texts:
                chunks.append(f"[slide {idx}] " + "\n".join(texts))
    except Exception as exc:
        chunks.append(f"[EXTRACT_ERROR] {exc}")
    return clean_text("\n".join(chunks))


def extract_all() -> dict[str, str]:
    TMP_DIR.mkdir(parents=True, exist_ok=True)
    extracted = {}
    for name in FILES:
        path = SOURCE_DIR / name
        suffix = path.suffix.lower()
        if suffix == ".pdf":
            text = extract_pdf(path)
        elif suffix == ".docx":
            text = extract_docx(path)
        elif suffix == ".pptx":
            text = extract_pptx(path)
        else:
            continue
        extracted[name] = text
    with NOTES_TXT.open("w", encoding="utf-8") as f:
        for name, text in extracted.items():
            f.write(f"\n\n===== {name} =====\n")
            f.write(text[:22000])
    return extracted


def keywords(text: str, limit: int = 28) -> list[str]:
    tokens = re.findall(r"[A-Za-z][A-Za-z\-]{3,}|[\u4e00-\u9fff]{2,8}", text.lower())
    counts = Counter(t for t in tokens if t not in STOPWORDS and not t.isdigit())
    return [w for w, _ in counts.most_common(limit)]


def source_stats(extracted: dict[str, str]) -> str:
    lines = []
    for name, text in extracted.items():
        lines.append(f"{name}: {len(text):,} chars | " + ", ".join(keywords(text, 12)))
    return "\n".join(lines)


def para(style, text):
    return Paragraph(text, style)


def make_styles():
    registerFont(UnicodeCIDFont("STSong-Light"))
    base = getSampleStyleSheet()
    styles = {}
    styles["tiny"] = ParagraphStyle(
        "tiny",
        parent=base["BodyText"],
        fontName="STSong-Light",
        fontSize=4.55,
        leading=5.05,
        spaceAfter=0.55,
        alignment=TA_LEFT,
        wordWrap="CJK",
        allowWidows=0,
        allowOrphans=0,
    )
    styles["micro"] = ParagraphStyle(
        "micro",
        parent=styles["tiny"],
        fontSize=4.2,
        leading=4.65,
        spaceAfter=0.35,
    )
    styles["head"] = ParagraphStyle(
        "head",
        parent=styles["tiny"],
        fontSize=6.2,
        leading=6.6,
        textColor=colors.white,
        backColor=colors.HexColor("#18324a"),
        spaceBefore=1.2,
        spaceAfter=1.2,
    )
    styles["subhead"] = ParagraphStyle(
        "subhead",
        parent=styles["tiny"],
        fontSize=5.25,
        leading=5.55,
        textColor=colors.HexColor("#0f3c5c"),
        spaceBefore=0.6,
        spaceAfter=0.3,
    )
    styles["title"] = ParagraphStyle(
        "title",
        parent=styles["tiny"],
        fontSize=8.5,
        leading=9,
        textColor=colors.HexColor("#102131"),
        spaceAfter=2,
    )
    return styles


def b(s: str) -> str:
    return f"<b>{s}</b>"


def section(title: str, items: list[str], styles) -> list:
    flow = [para(styles["head"], f"&nbsp;{title}")]
    for item in items:
        if item.startswith("## "):
            flow.append(para(styles["subhead"], item[3:]))
        else:
            flow.append(para(styles["tiny"], item))
    return flow


PAGE1 = [
    (
        "CRISPR / 基因编辑 / 贺建奎案",
        [
            f"{b('技术链条')}：CRISPR-Cas9 = guide RNA 定位 + Cas9 切割 + NHEJ/ HDR 修复；优点：便宜、快、可编程；风险：off-target、mosaicism、pleiotropy、多基因性状不可控、代际不可逆。",
            f"{b('体细胞 vs 生殖系')}：体细胞编辑只影响患者本人，伦理上更接近治疗；胚胎/生殖系编辑会传给后代，牵涉未来人的同意、物种层面的风险、优生学滑坡。",
            f"{b('治疗 vs 增强')}：治疗=恢复基本健康/减少疾病负担；增强=提升正常能力或偏好性状。灰区：HIV 抵抗、智力、身高、抗病能力。",
            f"{b('贺建奎案例')}：2018 宣称 CCR5 编辑双胞胎 Lulu/Nana，目标是 HIV 免疫；核心问题：医学必要性弱、替代方案存在、知情同意不足、伦理审查失范、风险未充分评估、科学共同体共识被绕开。",
            f"{b('法律后果')}：2019 中国法院以非法行医等判刑三年；案件推动中国强化人类遗传资源、医学伦理审查、胚胎基因编辑监管。",
            f"{b('正方论点')}：若技术成熟，可预防严重遗传病；给家庭生育自主权；减少疾病痛苦与医疗成本；科学探索推动治疗突破。",
            f"{b('反方论点')}：不可逆代际风险；未来儿童无法同意；社会不平等会转化为生物不平等；从治疗滑向定制婴儿；监管滞后于商业动机。",
            f"{b('答题模板')}：先区分 somatic/germline，再用 four principles（beneficence, non-maleficence, autonomy, justice）评估，最后落到 governance: moratorium? registry? ethics review? public deliberation?",
            f"{b('关键词')}：CCR5, off-target, mosaicism, informed consent, clinical necessity, reproductive autonomy, eugenics, precautionary principle, responsible innovation.",
        ],
    ),
    (
        "治疗性克隆 / 干细胞 / 生命伦理法",
        [
            f"{b('概念')}：therapeutic cloning = somatic cell nuclear transfer (SCNT)：取体细胞核 + 去核卵细胞 -> 囊胚 -> embryonic stem cells，用于再生医学；区别于 reproductive cloning（产生完整个体）。",
            f"{b('医学潜力')}：患者特异性细胞/组织、免疫排斥低、神经损伤/帕金森/糖尿病/脊髓损伤治疗模型、疾病机制研究、药物筛选。",
            f"{b('伦理冲突')}：胚胎道德地位；工具化生命；卵子来源与女性身体商品化；资源分配；从治疗性克隆滑向生殖性克隆。",
            f"{b('Christopher Reeve')}：以脊髓损伤患者身份支持胚胎干细胞研究，强调研究带来希望，不应因抽象恐惧阻断真实病痛中的人。",
            f"{b('英国 HFE Act 1990')}：监管胚胎研究、辅助生殖；常见边界：14-day rule（原条出现前）、许可制度、禁止不受控研究。",
            f"{b('Therapeutic Cloning Act 2001')}：围绕是否允许治疗性克隆展开；可写成“允许有限研究 + 严禁生殖性克隆 + 严格许可和审查”。",
            f"{b('伦理法庭结构')}：原告=胚胎权利/自然秩序/滑坡风险；被告=患者利益/科学自由/严格监管；法官平衡 harm, benefit, dignity, justice。",
        ],
    ),
    (
        "AI / LLM / Agent / Constitutional AI",
        [
            f"{b('LLM 基础')}：大语言模型通过预训练学习 next-token prediction，涌现出 summarization, translation, reasoning-like behavior；局限：hallucination、bias、不可解释、上下文窗口、数据污染。",
            f"{b('从 LLM 到 Agent')}：LLM + tools + memory + planning + feedback loop = AI agent；典型循环：observe -> plan -> act/tool use -> reflect -> revise；风险：越权调用、目标漂移、隐私泄露、自动化错误放大。",
            f"{b('Alignment')}：让系统行为符合人类意图与价值。技术路径：SFT, RLHF, RLAIF, red-teaming, constitutional rules, safety filters, evals。",
            f"{b('Constitutional AI')}：Anthropic 方法：用一组“宪法原则”指导模型自我批评和修订，再用 AI feedback 训练 harmless/helpful 模型；重点是把价值约束显式化、可审计化。",
            f"{b('Helpful-Honest-Harmless')}：有用=解决任务；诚实=承认不确定/不编造；无害=避免伤害、违法、歧视、隐私侵犯。",
            f"{b('AI will kill love?')}：可从三层答：亲密关系被替代/商业化/数据化；陪伴 AI 提供低成本情感支持；真正风险不是“爱消失”，而是依恋被平台设计和不对称知识操控。",
            f"{b('课堂讨论句')}：AI 不只是工具，而是重新组织 attention, agency, labor, intimacy 的基础设施。",
        ],
    ),
    (
        "NFC / RFID / 非接触支付",
        [
            f"{b('技术')}：NFC 基于 RFID，高频 13.56MHz，近距离通信；模式：读写、点对点、卡模拟；支付常用 tokenization + secure element/TEE + 动态密码。",
            f"{b('优势')}：速度快、低摩擦、可嵌入手机/卡/穿戴设备；提升交易效率，推动无现金社会和数字身份。",
            f"{b('风险')}：窃听、relay attack、中间人、恶意标签、丢失手机、后台过度采集、位置/消费画像、老年人和弱势群体被技术排除。",
            f"{b('防范')}：距离控制、加密认证、动态令牌、交易限额、生物识别/多因素、关闭不必要 NFC、系统更新、隐私最小化、监管透明。",
            f"{b('人文角度')}：便利不是中立价值；支付技术改变信任结构：从面对面信用转向平台/银行/设备的基础设施信任。",
        ],
    ),
    (
        "社会信用系统 / 数据治理",
        [
            f"{b('定义')}：社会信用不是单一分数，而是多部门信用信息、惩戒/激励名单、平台评分、地方治理试验的组合。",
            f"{b('治理逻辑')}：通过数据化记录、分类、预测和奖惩，把道德/法律/市场行为转成可计算信用；目标是提升 compliance 与治理效率。",
            f"{b('争议')}：透明度不足、申诉机制弱、数据错误、过度惩戒、隐私侵入、算法歧视、行为寒蝉效应。",
            f"{b('对比框架')}：西方 credit score 偏金融风险；中国社会信用更连接行政治理和公共秩序；平台评分则以商业注意力和消费行为为核心。",
            f"{b('答题句')}：关键不是“有没有分数”，而是谁收集数据、按什么规则计算、如何纠错、惩戒是否相称、个人能否退出。",
        ],
    ),
]


PAGE2 = [
    (
        "塑料污染 / 环境技术伦理",
        [
            f"{b('Plastic')}：优点=轻质、便宜、耐腐蚀、医疗/包装/运输价值高；问题=一次性消费、微塑料、生物累积、海洋污染、焚烧/填埋外部性。",
            f"{b('生命周期')}：原料开采 -> 生产 -> 使用 -> 回收/焚烧/填埋/泄漏；只谈“回收”会遮蔽上游生产和消费结构。",
            f"{b('治理')}：reduce/reuse/recycle、生产者责任延伸(EPR)、押金制、可降解材料但需警惕 greenwashing、源头减量优先。",
            f"{b('技术承诺')}：焚烧/净化/智能分类等方案能降低污染，但也可能制造“技术会替我们善后”的道德风险。",
            f"{b('可背句')}：环境问题不是单纯技术失败，而是便利、消费、市场激励和治理边界共同塑造的后果。",
        ],
    ),
    (
        "时间、媒介与现代性",
        [
            f"{b('Clock: Ribbon of Time')}：钟表把流动经验切割成可测量、可交易、可管理的单位；从自然时间/事件时间转向机械时间。",
            f"{b('社会影响')}：工业纪律、准时性、工资劳动、铁路/城市同步、学校和工厂的时间表；时间成为治理身体和组织生产的工具。",
            f"{b('人文问题')}：效率提升同时带来异化：人按机器节奏生活；“拥有时间”变成现代自由与焦虑的核心。",
            f"{b('答题连接')}：可连到 Foucault 的规训权力：钟表/时间表/考勤把身体纳入可观察、可比较、可优化的秩序。",
        ],
    ),
    (
        "福柯 / 阿甘本 / 生命政治",
        [
            f"{b('Foucault')}：biopolitics = 现代权力从“让人死/让人活”转向“使人活并管理人口”：出生率、死亡率、卫生、风险、疾病、统计。",
            f"{b('Discipline vs Biopower')}：规训针对个体身体（学校、军队、医院、工厂）；生命权力针对人口（公共卫生、保险、城市治理、算法风险）。",
            f"{b('Agamben')}：bare life（赤裸生命）与 state of exception（例外状态）；主权通过决定谁被纳入法律保护、谁被排除为可牺牲生命而显现。",
            f"{b('差异')}：福柯强调治理术、人口管理、微观权力网络；阿甘本强调主权决断、集中营范式、法律内外的边界。",
            f"{b('内在联结')}：二者都说明现代政治把“生命本身”纳入权力对象；科技治理、医学伦理、疫情管控、基因编辑、数据评分都可用生命政治分析。",
            f"{b('万能句')}：当技术把生命转化为数据、风险和可优化对象时，治理效率与主体自由之间的张力就出现了。",
        ],
    ),
    (
        "文学与技术批判：Hard Times",
        [
            f"{b('Dickens《艰难时世》')}：批判功利主义、工业资本主义和“Facts only”的教育观；Coketown 象征机器化城市，烟囱、重复、灰色秩序压扁想象力。",
            f"{b('人物功能')}：Gradgrind=事实/计算理性；Bounderby=资本自我神话；Sissy=同情、想象、经验知识；Louisa=被功利教育伤害的主体。",
            f"{b('课程连接')}：科技理性若脱离情感、伦理和叙事，会把人变成指标/劳动力/数据点；人文不是反科学，而是问科学服务何种生活。",
        ],
    ),
    (
        "余数生命悖论 / 医学与价值排序",
        [
            f"{b('问题意识')}：“剩余生命”把人按可预期寿命、生产力、治疗收益排序，容易在资源分配中制造隐性不平等。",
            f"{b('伦理张力')}：最大化总生命年(QALY/效益主义) vs 每个人平等尊严(义务论/权利论) vs 优先弱者/最差者优先。",
            f"{b('课堂应用')}：器官移植、ICU 分流、罕见病药物、高价疗法、老龄化社会。问：谁定义“值得救”？指标是否遮蔽亲情、痛苦、社会关系？",
        ],
    ),
    (
        "课堂展示快速答题框架",
        [
            f"{b('三层分析')}：1 技术原理/功能；2 社会后果/利益相关者；3 伦理规范/治理方案。",
            f"{b('四原则')}：autonomy 自主/知情同意；beneficence 促进福祉；non-maleficence 不伤害；justice 公平分配。",
            f"{b('利益相关者')}：研发者、企业、政府、用户/患者、未来世代、弱势群体、环境/非人生命。",
            f"{b('风险类型')}：安全风险、隐私风险、歧视风险、不可逆风险、责任归属、技术锁定、数字鸿沟、商业操控。",
            f"{b('治理工具')}：许可/审查、透明解释、第三方审计、数据最小化、公众参与、责任追踪、红线禁区、沙盒试验、国际共识。",
            f"{b('万能开头')}：This case shows that technology is never merely technical; it redistributes power, risk, responsibility and imagination.",
            f"{b('万能结尾')}：A reasonable position is neither blind optimism nor rejection, but conditional acceptance under proportional risk control, public deliberation and accountable governance.",
            f"{b('中英关键词')}：responsible innovation 负责任创新; precautionary principle 预防原则; informed consent 知情同意; human dignity 人的尊严; distributive justice 分配正义; surveillance 监控; autonomy 自主性; governance 治理。",
        ],
    ),
]


PAGE1_EXTRA = [
    (
        "CRISPR 临床细节速记",
        [
            f"{b('ex vivo')}：取出细胞 -> 体外编辑/筛选/扩增 -> 回输；优点是可检测 edited cells，降低体内递送难题；缺点是步骤复杂、成本高、适合血液/免疫细胞。",
            f"{b('in vivo')}：编辑工具直接进入体内目标组织；优点是可治疗体内器官疾病；难点是 delivery、免疫反应、不可控扩散、长期随访。",
            f"{b('疾病例子')}：HIV: CCR5/CXCR4 knockout；癌症: PD-1 knockout T cells, CAR-T/TRAC locus；血液病: β-globin/HBB correction, sickle cell disease, CD34+ HSPCs。",
            f"{b('技术比较')}：ZFN/TALEN 需要蛋白工程；CRISPR 依靠 guide RNA 重新编程，设计成本低；base editing/prime editing 试图减少双链断裂。",
            f"{b('实验指标')}：on-target efficiency, off-target rate, indel profile, HDR/NHEJ ratio, translocation, cell viability, engraftment, immune response。",
        ],
    ),
    (
        "贺建奎/监管句库",
        [
            f"{b('问题定性')}：不是“科学突破 vs 保守伦理”的简单冲突，而是临床必要性、风险比例、审查程序、商业/名誉动机与未来儿童权利的共同失败。",
            f"{b('可写批判')}：The experiment transformed babies into a proof-of-concept platform before safety, necessity and consent were established.",
            f"{b('治理建议')}：国际登记、胚胎编辑红线、伦理审查独立性、长期随访、透明发表、利益冲突披露、公共参与。",
        ],
    ),
    (
        "AI 高频风险与对策",
        [
            f"{b('幻觉')}：模型生成流畅但不真实的内容；应使用 retrieval、引用、verification、uncertainty disclosure。",
            f"{b('偏见')}：训练数据反映社会偏见；应做 bias eval、数据治理、人工复核、影响评估。",
            f"{b('代理风险')}：agent 会连续行动，错误可累积；应设置权限边界、approval gate、日志审计、回滚机制。",
            f"{b('亲密关系')}：AI companion 的问题不是“有没有情感”，而是信息不对称、依赖设计、商业化陪伴和人格模拟是否操纵用户。",
        ],
    ),
    (
        "NFC 诈骗三阶段",
        [
            f"{b('1 身份伪造')}：冒充客服/航空/银行，制造退款、改签、保险返还等紧迫情境。",
            f"{b('2 恶意引导')}：要求下载 App、开屏幕共享、贴近银行卡/手机 NFC 区域，读取卡号、有效期、CVV 或绑定支付。",
            f"{b('3 无感转账')}：利用免密/默认支付/远控完成盗刷；用户常在大额交易后才发现异常。",
            f"{b('防线')}：官方渠道核验；不装陌生 App；关闭默认支付/免密；降低限额；敏感操作用多因素；银行卡远离手机背部 NFC 区。",
        ],
    ),
]


PAGE2_EXTRA = [
    (
        "展示/问答万能结构",
        [
            f"{b('定义题')}：X is not merely a device/application; it is an infrastructure that reorganizes relations among people, institutions and values.",
            f"{b('案例题')}：fact -> stakeholder -> benefit -> harm -> principle -> governance。每一步一句，不要只表态。",
            f"{b('比较题')}：先列共同点，再抓核心差异：对象不同、权力机制不同、风险尺度不同、治理工具不同。",
            f"{b('反驳句')}：This objection is valid, but it assumes technology is socially neutral. In practice, design choices distribute risk unevenly.",
            f"{b('平衡句')}：The question is not whether to ban or embrace it, but under what conditions it can be justified.",
        ],
    ),
    (
        "常用理论钩子",
        [
            f"{b('Technological determinism')}：技术决定论，把社会变化过度归因于技术本身。",
            f"{b('Social construction of technology')}：技术由利益、制度、文化共同塑造。",
            f"{b('Risk society')}：现代社会的核心不只是生产财富，也生产和分配风险。",
            f"{b('Surveillance capitalism')}：数据收集、预测和行为影响构成新的商业权力。",
            f"{b('Digital divide')}：技术收益与风险按阶层、年龄、地区、能力不均分布。",
        ],
    ),
    (
        "可直接套用的结论",
        [
            "科技扩展人的能力，也重写人的边界：谁被看见、谁被分类、谁承担风险、谁获得收益。",
            "人文视角不是反技术，而是要求技术接受价值追问：为了谁、由谁决定、代价给谁、能否纠错。",
            "好的治理不只限制坏行为，也要让普通人理解规则、参与规则，并在受损时获得救济。",
        ],
    ),
]


PAGE1_MORE = [
    (
        "CRISPR 展示细节",
        [
            f"{b('Victoria Gray')}：镰状细胞病患者案例常用来说明基因疗法的治疗价值：从长期疼痛/住院到被视为“重生”的医学希望。",
            f"{b('deaf culture')}：耳聋不只可被定义为缺陷，也可能是一种语言与文化身份；基因编辑会遇到“谁有权定义正常”的问题。",
            f"{b('隐喻陷阱')}：把基因说成“代码/说明书/命运”会夸大可控性；真实生命是基因、环境、社会条件共同作用。",
            f"{b('定义霸权')}：当医学、市场和国家把某种差异定义成 fault，父母可能承受“生下不完美孩子是道德失败”的压力。",
            f"{b('公平问题')}：若癌症免疫/遗传病疗法价格数百万，问题会从“能不能做”变为“谁能活得更好”。",
        ],
    ),
    (
        "克隆伦理法庭攻防",
        [
            f"{b('支持许可')}：阻止研究并非中立；inaction also has victims。患者正在承受可避免的痛苦，严格监管下研究可被道德化。",
            f"{b('反对许可')}：胚胎不是普通材料；把早期生命当作可制造、可消耗的工具，会改变社会理解 human life 的方式。",
            f"{b('中间立场')}：承认胚胎具有某种 moral status，但不等于完整人格；可允许剩余胚胎/SCNT 研究，同时禁止生殖性克隆。",
            f"{b('关键追问')}：卵子来源是否剥削？14 天规则是否充分？商业专利会不会让患者希望被垄断？",
        ],
    ),
    (
        "AI Agent / Medvi / OPC",
        [
            f"{b('Medvi prompt')}：tiny team + AI Agents build fast-growing company：是创业民主化，还是监管风险，或两者兼有？",
            f"{b('能力替代')}：agents 可替代/增强客服、销售线索、市场文案、代码原型、数据分析、流程自动化。",
            f"{b('创始人角色变化')}：从亲自执行转向 orchestration：定义目标、拆任务、审输出、管风险、承担责任。",
            f"{b('民主化')}：降低启动成本、小团队获得大公司能力、提高边缘创业者机会。",
            f"{b('风险')}：责任主体模糊、合规滞后、幻觉造成商业损失、数据泄露、劳动替代、监管套利。",
        ],
    ),
    (
        "NFC 消费主义补丁",
        [
            f"{b('一触即付')}：降低支付摩擦，弱化“花钱痛感”；心理账户更容易被绕开，消费从生存型转向享乐型。",
            f"{b('课堂句')}：Quick payment changes not only how we pay, but how we desire, calculate and justify consumption.",
            f"{b('风险权衡')}：0.01% 小概率盗刷 vs 高频便利；可用 expected loss + trust + habit 分析。",
            f"{b('身份门禁')}：access card cloning 的损害不是钱，而是空间安全、身体安全和组织信任。",
        ],
    ),
    (
        "社会信用细节",
        [
            f"{b('数据来源')}：金融、司法、行政处罚、市场监管、平台行为、交通/公共服务等多维记录。",
            f"{b('机制')}：红名单/黑名单、联合惩戒、信用修复、跨部门数据共享；重点是 compliance infrastructure。",
            f"{b('风险')}：unfair scoring、数据过期/错误、算法黑箱、惩戒过度、私人平台与国家治理边界模糊。",
            f"{b('可答')}：信用治理若要正当，必须满足合法性、透明性、比例原则、申诉纠错、最小必要数据。",
        ],
    ),
]


PAGE2_MORE = [
    (
        "塑料 PPT 细节",
        [
            f"{b('传统方案三缺陷')}：填埋降解需数百年且污染土壤地下水；人工分类贵且低效；回收质量不稳定，难形成可持续闭环。",
            f"{b('中国垃圾难烧')}：厨余/塑料混合、含氯塑料(PVC)可能带来二噁英、腐蚀和复杂净化问题。",
            f"{b('技术方案')}：炉排、耐腐蚀合金、多污染物协同净化等；PPT 强调净化效率与现实垃圾山案例。",
            f"{b('核心矛盾')}：技术越能“无害化”，越可能削弱源头减塑意识。",
            f"{b('Marx/中立性')}：工具本身不决定善恶；社会关系、所有权、激励和使用目的决定技术后果。",
        ],
    ),
    (
        "钟表 PPT 细节",
        [
            f"{b('正效应')}：同步协作、航海/铁路/工业生产、科学测量、公共生活秩序。",
            f"{b('负效应')}：time anxiety、感知钝化、压力疲劳、过长工作时间；人变成按时刻表运转的机器。",
            f"{b('Modern Stephen')}：可联系《艰难时世》中被工业秩序压迫的劳动者，现代版本是被 KPI、排班、通知和算法节奏驱动的人。",
            f"{b('Thoreau')}：过度劳动让人“没有时间成为人”；用于反驳单一效率崇拜。",
        ],
    ),
    (
        "铁路材料：技术背后的人",
        [
            f"{b('系统分工')}：线路规划/技术研发、施工建设、调度指挥、司机、乘务、维修、货运共同构成铁路技术系统。",
            f"{b('技术不是机器本身')}：安全运行依赖制度、训练、劳动纪律、现场判断和跨岗位协作。",
            f"{b('人文连接')}：高铁的便利背后有不可见劳动、家庭缺席、职业风险和地方文化差异。",
            f"{b('可用句')}：Infrastructure hides labor: when the system works smoothly, the workers who keep it safe become invisible.",
        ],
    ),
    (
        "AI Will Kill Love 辩论素材",
        [
            f"{b('正方')}：AI 提供 zero-cost, conflict-free, always-obedient intimacy，削弱真实关系所需的 patience, compromise, responsibility。",
            f"{b('反方')}：AI 可作为镜子/陪伴/心理支持，帮助人练习表达、修复孤独，而不是替代现实关系。",
            f"{b('关键分歧')}：love 的本质是情绪满足，还是对他者的责任？AI 的“服从”能否构成 reciprocal relationship？",
            f"{b('好句')}：Real love is not frictionless comfort; it is the willingness to stay with uncertainty and another person's freedom.",
        ],
    ),
    (
        "AI Safety deck 关键词",
        [
            f"{b('Adversarial')}：微小扰动/恶意输入可让模型输出失败；安全不是平均表现，而是边界条件表现。",
            f"{b('Hallucination')}：信心不等于事实；mitigation = retrieval, citation, tool verification, refusal when uncertain。",
            f"{b('Bias')}：产品、制度、职业、历史数据都可能固化刻板印象。",
            f"{b('Causality')}：相关性不等于因果；action gap 指模型建议到现实行动之间还有制度和责任鸿沟。",
            f"{b('Alignment ambiguity')}：人类价值本身冲突，trolley dilemma 展示 utilitarianism vs absolute morality。",
        ],
    ),
]


PAGE1_GLOSS = [
    (
        "基因编辑术语表",
        [
            "gRNA/sgRNA: 引导 Cas9 到目标序列；PAM: Cas9 识别所需短序列；DSB: double-strand break；NHEJ: 易产生 indel；HDR: 依赖模板精确修复。",
            "off-target: 非目标位点编辑；mosaicism: 胚胎细胞编辑结果不一致；germline: 可遗传；somatic: 不遗传；enhancement: 超越治疗。",
            "clinical necessity: 是否有真实医疗必要；risk-benefit ratio: 风险收益比；long-term follow-up: 长期随访；global moratorium: 国际暂停。",
            "eugenics: 优生学；designer baby: 定制婴儿；reproductive autonomy: 生育自主；future generations: 未来世代。",
        ],
    ),
    (
        "克隆/干细胞术语表",
        [
            "SCNT: somatic cell nuclear transfer；ntESC: nuclear-transfer embryonic stem cells；blastocyst: 囊胚；oocyte: 卵母细胞。",
            "moral status: 道德地位；potentiality argument: 潜能论；instrumentalization: 工具化；commodification: 商品化。",
            "14-day rule: 原条出现前研究边界；licensing: 许可制；HFEA: 英国人类受精与胚胎管理框架。",
            "支持方关键词：patient suffering, regenerative medicine, cell replacement, hope, strict oversight；反对方关键词：dignity, slippery slope, sanctity of life。",
        ],
    ),
    (
        "AI 治理术语表",
        [
            "SFT: supervised fine-tuning；RLHF: human feedback；RLAIF: AI feedback；red teaming: 对抗测试；evals: 评测。",
            "reward hacking: 钻奖励函数空子；goal misgeneralization: 目标泛化错误；overreliance: 过度依赖；automation bias: 自动化偏见。",
            "data leakage: 数据泄露；model inversion: 反推训练数据；prompt injection: 提示注入；tool abuse: 工具滥用。",
            "guardrail: 护栏；human-in-the-loop: 人在回路；audit trail: 审计链；least privilege: 最小权限。",
        ],
    ),
    (
        "支付/信用/数据术语表",
        [
            "tokenization: 用令牌替代真实卡号；secure element/TEE: 安全执行环境；relay attack: 中继攻击；skimming: 近距离盗刷。",
            "mental accounting: 心理账户；frictionless payment: 无摩擦支付；consumerism: 消费主义；hedonic consumption: 享乐消费。",
            "blacklist/whitelist: 黑/红名单；joint punishment: 联合惩戒；credit repair: 信用修复；proportionality: 比例原则。",
            "data minimization: 数据最小化；purpose limitation: 目的限定；due process: 正当程序；right to explanation: 解释权。",
        ],
    ),
]


PAGE2_GLOSS = [
    (
        "环境/时间/文学术语表",
        [
            "externality: 外部性；EPR: extended producer responsibility；NIMBY: 邻避效应；greenwashing: 绿色洗白；moral hazard: 道德风险。",
            "mechanical time: 机械时间；time discipline: 时间规训；alienation: 异化；instrumental reason: 工具理性。",
            "utilitarianism: 功利主义；Facts only: 《艰难时世》里被批判的狭窄事实教育；imagination: 人文想象力。",
            "efficiency trap: 效率陷阱；labor invisibility: 劳动不可见；infrastructure: 基础设施。",
        ],
    ),
    (
        "生命政治术语表",
        [
            "biopower: 生命权力；anatomo-politics: 身体规训；population: 人口；governmentality: 治理术。",
            "sovereignty: 主权；state of exception: 例外状态；bare life: 赤裸生命；inclusive exclusion: 排除性纳入。",
            "apparatus/dispositif: 装置；normalization: 正常化；security dispositif: 安全装置；pastoral power: 牧领权力。",
            "Foucault = 微观治理/人口/现代性；Agamben = 主权/法律/例外/西方政治结构。",
        ],
    ),
    (
        "答题短句库",
        [
            "Technology is a social choice disguised as a technical solution.",
            "Convenience often works by hiding labor, risk and responsibility.",
            "The ethical question begins when the person affected cannot meaningfully consent.",
            "A tool becomes political when it classifies, excludes or disciplines people.",
            "Safety should be measured not only by average performance, but by failure modes.",
            "Public trust requires transparency, appeal mechanisms and proportional sanctions.",
        ],
    ),
    (
        "压轴万能句",
        [
            "科技与人文的核心张力：能力扩张 vs 价值约束；效率提升 vs 尊严/自由；个体便利 vs 集体风险；当代收益 vs 未来代价。",
            "写作顺序：先承认技术价值，再指出风险分配不均，最后给出条件性接受：limited, transparent, accountable, reversible where possible。",
            "不要只说“有利有弊”；要说“什么条件下利大于弊，谁来判断，出了问题谁负责”。",
        ],
    ),
]


PAGE1_COLUMNS = [
    [
        (
            "CRISPR / 基因编辑",
            [
                f"{b('CRISPR-Cas9')}：细菌适应性免疫系统改造而来；sgRNA 定位目标 DNA，Cas9 造成双链断裂，细胞通过 NHEJ 或 HDR 修复。",
                f"{b('关键部件')}：crRNA/tracrRNA 或 sgRNA；PAM 序列决定 Cas9 能否结合；HNH/RuvC 结构域切割双链；dCas9 可做转录调控而非切割。",
                f"{b('NHEJ/HDR')}：NHEJ 快但易产生 indel，适合 knockout；HDR 依赖模板，可精确替换/修复但效率低，细胞周期限制明显。",
                f"{b('ex vivo')}：取出细胞，体外编辑、筛选、扩增、回输；适合血液/免疫细胞，优点是可检测 edited cells，缺点是成本和流程复杂。",
                f"{b('in vivo')}：编辑工具直接进入体内目标组织；难点是递送、免疫反应、组织特异性、长期脱靶和不可逆后果。",
                f"{b('疾病应用')}：HIV: CCR5/CXCR4 knockout；肿瘤: PD-1 knockout T cells, CAR-T/TRAC locus；血液病: HBB/β-globin correction, CD34+ HSPCs。",
                f"{b('Victoria Gray')}：镰状细胞病基因治疗案例，展示从长期疼痛/住院到“医学重生”的治疗价值，也凸显高价疗法的公平性问题。",
                f"{b('实验指标')}：on-target efficiency, off-target rate, indel profile, HDR/NHEJ ratio, translocation, cell viability, engraftment, immune response。",
                f"{b('技术演进')}：ZFN/TALEN 依赖蛋白工程；CRISPR 通过 guide RNA 重新编程；base editing/prime editing 试图减少双链断裂。",
                f"{b('核心风险')}：off-target、mosaicism、pleiotropy、多基因性状不可控、长期随访不足、代际不可逆、治疗/增强边界漂移。",
            ],
        ),
        (
            "贺建奎案",
            [
                f"{b('事件')}：2018 年宣称 CCR5 编辑双胞胎 Lulu/Nana，目标是 HIV 免疫；2019 年中国法院判刑三年。",
                f"{b('科学问题')}：CCR5 与 HIV 感染相关，但 HIV 有成熟阻断方案；CCR5 还与免疫/神经等复杂功能有关，简单 knockout 并非无代价。",
                f"{b('伦理问题')}：医学必要性弱、替代方案存在、知情同意不足、伦理审查失范、未来儿童无法同意、风险由后代承担。",
                f"{b('社会问题')}：从治疗滑向增强/定制婴儿；社会不平等转化为生物不平等；父母可能承受“不优化孩子就是失责”的压力。",
                f"{b('监管后果')}：推动中国强化人类遗传资源、医学伦理审查、胚胎基因编辑红线、长期随访与科研责任追踪。",
            ],
        ),
        (
            "基因编辑术语",
            [
                "somatic editing: 只影响患者本人；germline editing: 可遗传给后代；enhancement: 超越治疗目标。",
                "clinical necessity: 临床必要性；risk-benefit ratio: 风险收益比；reproductive autonomy: 生育自主；eugenics: 优生学。",
                "deaf culture: 耳聋可被理解为语言/文化身份，不只是缺陷；基因编辑触及“谁定义正常”的权力。",
                "designer baby: 定制婴儿；future generations: 未来世代；global moratorium: 国际暂停或禁令倡议。",
            ],
        ),
    ],
    [
        (
            "治疗性克隆 / 干细胞",
            [
                f"{b('SCNT')}：somatic cell nuclear transfer，体细胞核 + 去核卵母细胞 -> 囊胚 -> embryonic stem cells；目标是再生医学，不是复制完整个体。",
                f"{b('therapeutic vs reproductive')}：治疗性克隆用于细胞/组织/疾病模型；生殖性克隆产生完整个体，伦理和法律禁止程度更高。",
                f"{b('医学潜力')}：患者特异性细胞、免疫排斥低、神经损伤/帕金森/糖尿病/脊髓损伤模型、药物筛选、疾病机制研究。",
                f"{b('Christopher Reeve')}：以脊髓损伤患者身份支持胚胎干细胞研究，强调真实患者痛苦与研究希望之间的道德重量。",
                f"{b('主要争议')}：胚胎道德地位、潜能论、生命工具化、卵子来源与女性身体商品化、商业专利、资源分配、公私利益冲突。",
                f"{b('14-day rule')}：原条出现前作为胚胎研究边界；试图在早期生命道德地位与科研价值之间建立制度折中。",
                f"{b('HFE Act 1990')}：英国框架监管辅助生殖和胚胎研究，强调许可制、审查、信息记录和禁止不受控研究。",
                f"{b('Therapeutic Cloning Act 2001')}：围绕治疗性克隆许可与生殖性克隆禁令展开，典型结构是有限研究 + 严格许可 + 禁止生殖用途。",
                f"{b('卵子问题')}：oocyte availability 既是技术瓶颈，也是伦理问题，涉及捐赠、补偿、健康风险和身体剥削。",
                f"{b('术语')}：ntESC, blastocyst, oocyte, potentiality argument, instrumentalization, commodification, sanctity of life。",
            ],
        ),
        (
            "AI / LLM / Agent",
            [
                f"{b('LLM')}：通过大规模预训练学习 next-token prediction，表现出 summarization, translation, reasoning-like behavior；不是事实数据库。",
                f"{b('局限')}：hallucination、bias、上下文窗口、不可解释、训练数据污染、信心与正确性脱钩。",
                f"{b('Agent')}：LLM + tools + memory + planning + feedback loop；循环为 observe -> plan -> act/tool use -> reflect -> revise。",
                f"{b('Agent 风险')}：越权调用、目标漂移、数据泄露、工具滥用、自动化错误放大、责任主体模糊。",
                f"{b('Medvi/OPC')}：AI Agents 让 tiny team 获得客服、销售线索、市场文案、代码原型、数据分析和流程自动化能力；创始人从执行者转为 orchestrator。",
                f"{b('创业双面性')}：降低启动成本和组织门槛，同时带来合规滞后、监管套利、劳动替代、幻觉导致商业损失。",
            ],
        ),
    ],
    [
        (
            "Constitutional AI / AI Safety",
            [
                f"{b('Alignment')}：让系统行为符合人类意图和价值；技术路径包括 SFT, RLHF, RLAIF, red-teaming, evals, safety filters。",
                f"{b('Constitutional AI')}：Anthropic 用一组宪法原则指导模型自我批评和修订，再用 AI feedback 训练 helpful/harmless 模型。",
                f"{b('HHH')}：Helpful = 解决任务；Honest = 承认不确定、不编造；Harmless = 避免伤害、违法、歧视、隐私侵犯。",
                f"{b('reward hacking')}：模型钻奖励函数空子；goal misgeneralization：训练目标在新情境中偏离；overreliance：人过度相信自动化结果。",
                f"{b('Adversarial')}：微小扰动/恶意输入可让模型失败；安全不是平均表现，而是边界条件和失败模式。",
                f"{b('Hallucination')}：流畅输出不等于事实；缓解依赖 retrieval、citation、tool verification、uncertainty disclosure。",
                f"{b('Bias')}：产品、制度、职业、历史数据会固化刻板印象；需要数据治理、评估、复核和影响审查。",
                f"{b('Causality')}：相关性不等于因果；action gap 指模型建议到现实行动之间仍有制度、责任和执行鸿沟。",
                f"{b('trolley dilemma')}：展示 utilitarianism 与 absolute morality 的冲突；alignment 难点在于人类价值本身不统一。",
            ],
        ),
        (
            "AI 与爱",
            [
                f"{b('核心争点')}：AI 是否削弱真实亲密关系所需的 patience, compromise, responsibility, uncertainty。",
                f"{b('替代风险')}：AI companion 提供 24/7、低成本、低冲突、顺从式陪伴，使用户回避真实关系中的摩擦和他者自由。",
                f"{b('支持价值')}：AI 可作为镜子、陪伴和心理支持，帮助孤独者练习表达、稳定情绪、进入现实关系。",
                f"{b('平台问题')}：真正风险不只是“爱消失”，而是依恋被商业设计、数据画像和不对称知识操控。",
            ],
        ),
    ],
    [
        (
            "NFC / RFID / 支付",
            [
                f"{b('NFC')}：Near Field Communication，基于 RFID，13.56MHz，几厘米近距离通信；区别于 Bluetooth/Wi-Fi 的配对流程。",
                f"{b('三种模式')}：Reader/Writer 读写标签；Peer-to-Peer 两设备交换信息；Card Emulation 手机模拟交通卡、门禁卡或银行卡。",
                f"{b('支付安全')}：tokenization 替代真实卡号；secure element/TEE 保护密钥；动态密码和交易限额降低盗刷风险。",
                f"{b('便利')}：移动支付、公交地铁、门禁、学生卡、图书馆、智能标签、快速配对；提高交易效率和日常服务无缝性。",
                f"{b('风险')}：skimming、relay attack、中间人、恶意 App、屏幕共享远控、默认支付、免密额度、位置/消费画像。",
                f"{b('NFC 诈骗')}：身份伪造 -> 下载恶意 App/屏幕共享/贴卡读取 -> 绑定支付或无感转账；南京案例损失接近 12 万元。",
                f"{b('防范')}：关闭不必要 NFC；高风险交易多因素认证；降低免密额度；要求手机解锁；不装陌生 App；卡证远离手机背部 NFC 区。",
                f"{b('消费影响')}：一触即付降低支付摩擦，削弱花钱痛感；心理账户被绕开，消费从必要型转向享乐型。",
            ],
        ),
        (
            "社会信用 / 数据治理",
            [
                f"{b('定义')}：社会信用不是单一分数，而是信用信息、红黑名单、联合惩戒、信用修复、地方治理试验和平台评分的组合。",
                f"{b('数据来源')}：金融、司法、行政处罚、市场监管、交通、公共服务、平台行为等多维记录。",
                f"{b('治理逻辑')}：把行为数据化、分类、预测和奖惩，形成 compliance infrastructure，提高执行效率和公共秩序管理。",
                f"{b('争议')}：unfair scoring、数据错误/过期、算法黑箱、惩戒过度、隐私侵入、申诉机制弱、寒蝉效应。",
                f"{b('边界')}：西方 credit score 偏金融风险；中国社会信用连接行政治理；平台评分服务商业注意力和消费行为。",
                f"{b('正当性条件')}：合法性、透明性、比例原则、申诉纠错、数据最小必要、目的限定、解释权。",
            ],
        ),
    ],
]


PAGE2_COLUMNS = [
    [
        (
            "塑料污染 / 环境技术伦理",
            [
                f"{b('问题')}：塑料轻质、便宜、耐腐蚀，但一次性消费导致微塑料、生物累积、海洋污染、填埋/焚烧外部性。",
                f"{b('生命周期')}：原料开采 -> 生产 -> 使用 -> 回收/焚烧/填埋/泄漏；只谈回收会遮蔽上游生产和消费结构。",
                f"{b('传统方案缺陷')}：填埋降解需数百年并污染土壤地下水；人工分类贵且低效；回收质量不稳定，难形成闭环。",
                f"{b('中国垃圾难烧')}：厨余与塑料混合，PVC 等含氯塑料可能带来二噁英、腐蚀和复杂净化问题。",
                f"{b('技术方案')}：炉排、耐腐蚀合金、多污染物协同净化、智能分类；技术越能无害化，越可能削弱源头减塑意识。",
                f"{b('Moral hazard')}：当居民相信“技术会处理一切污染”，过度塑料消费会被合理化。",
                f"{b('治理')}：reduce/reuse/recycle、EPR、押金制、源头减量、消费结构调整、警惕 greenwashing。",
                f"{b('Marx 技术中立')}：工具本身不决定善恶；社会关系、所有权、激励和使用目的决定技术后果。",
            ],
        ),
        (
            "钟表 / 时间现代性",
            [
                f"{b('Clock')}：钟表把流动经验切割成可测量、可交易、可管理的单位，从自然时间/事件时间转向机械时间。",
                f"{b('正效应')}：航海、铁路、工业生产、科学测量、公共生活同步、学校与城市秩序。",
                f"{b('负效应')}：time anxiety、感知钝化、压力疲劳、过长工作时间；人按时刻表和机器节奏生活。",
                f"{b('time discipline')}：时间表、考勤、排班把身体纳入可观察、可比较、可优化的秩序。",
                f"{b('Modern Stephen')}：现代劳动者被 KPI、排班、通知和算法节奏驱动，延续工业秩序对人的压迫。",
                f"{b('Thoreau')}：过度劳动让人没有时间成为人；批判单一效率崇拜。",
            ],
        ),
    ],
    [
        (
            "福柯 / 阿甘本 / 生命政治",
            [
                f"{b('Foucault')}：biopolitics = 现代权力从“让人死/让人活”转向“使人活并管理人口”；对象是出生率、死亡率、卫生、风险、疾病、统计。",
                f"{b('discipline')}：规训针对个体身体，典型场景是学校、军队、医院、工厂、监狱。",
                f"{b('biopower')}：生命权力针对人口，典型机制是公共卫生、保险、城市治理、风险模型和安全装置。",
                f"{b('governmentality')}：治理术强调以自由主义、统计、专业知识和制度技术管理人口。",
                f"{b('Agamben')}：bare life 与 state of exception；主权通过决定谁被纳入法律保护、谁被排除为可牺牲生命而显现。",
                f"{b('inclusive exclusion')}：赤裸生命以被排除的方式被纳入政治秩序；例外状态显示法内/法外边界。",
                f"{b('差异')}：福柯重微观治理、人口管理、现代性；阿甘本重主权决断、法律边界、例外状态和西方政治结构。",
                f"{b('内在联结')}：阿甘本不是简单背离福柯，而是在主权-法律层面扩展福柯的生命政治问题域，并把主权内化于治理。",
                f"{b('装置')}：apparatus/dispositif 捕捉、引导、塑造和控制生命；技术系统、医学制度、数据平台都可成为装置。",
                f"{b('应用')}：疫情管控、基因编辑、医疗分流、社会信用、AI 风险评分都把生命转化为数据、风险和治理对象。",
            ],
        ),
        (
            "余数生命悖论 / 医学分配",
            [
                f"{b('问题')}：“剩余生命”把人按预期寿命、生产力、治疗收益排序，容易在资源分配中制造隐性不平等。",
                f"{b('QALY')}：质量调整生命年强调效益最大化；优点是可比较，问题是可能压低老人、残障者和慢性病患者价值。",
                f"{b('冲突')}：效益主义最大化总生命年；权利论强调每个人平等尊严；弱者优先强调最差者优先。",
                f"{b('场景')}：器官移植、ICU 分流、罕见病药物、高价基因疗法、老龄化社会、疫情资源紧缺。",
            ],
        ),
    ],
    [
        (
            "文学与技术批判：Hard Times",
            [
                f"{b('Dickens')}：《艰难时世》批判功利主义、工业资本主义和 Facts only 教育观。",
                f"{b('Coketown')}：机器化城市象征；烟囱、重复、灰色秩序和污染把生活压成生产节奏。",
                f"{b('Gradgrind')}：事实/计算理性，把儿童当作知识容器和可训练对象。",
                f"{b('Bounderby')}：资本自我神话，掩盖劳动剥削与阶级结构。",
                f"{b('Sissy')}：同情、想象、经验知识，代表被功利教育排斥的人文能力。",
                f"{b('Louisa')}：被事实教育伤害的主体，说明情感和想象贫乏会造成真实生活破裂。",
                f"{b('课程连接')}：科技理性若脱离情感、伦理和叙事，会把人变成指标、劳动力和数据点。",
            ],
        ),
        (
            "铁路 / 基础设施劳动",
            [
                f"{b('系统分工')}：线路规划、技术研发、施工建设、调度指挥、司机、乘务、维修、货运共同构成铁路技术系统。",
                f"{b('工程端')}：路线勘测、桥梁隧道、路基、轨道、信号、电力、制动系统共同支撑安全运行。",
                f"{b('运行端')}：调度中心 24 小时协调车次、线路、客流、天气和突发情况；司机承担全程驾驶和安全责任。",
                f"{b('维护端')}：巡检、抢修、设备保障、荒漠/山区/严寒线路维护，体现不可见劳动和职业风险。",
                f"{b('人文面')}：高铁便利背后有家庭缺席、身体消耗、地方文化差异、乘务沟通和安全仪式。",
                f"{b('Slides for pre')}：Yubari、Liangshan、Chicago 等地方铁路展示技术基础设施与地方经济、记忆和衰落之间的关系。",
            ],
        ),
    ],
    [
        (
            "科技与社会理论",
            [
                f"{b('technological determinism')}：技术决定论，过度把社会变化归因于技术本身。",
                f"{b('SCOT')}：social construction of technology，技术由利益、制度、文化、用户实践共同塑造。",
                f"{b('risk society')}：现代社会不仅生产财富，也生产和分配风险；风险常跨地域、跨世代、跨阶层。",
                f"{b('surveillance capitalism')}：数据收集、预测和行为影响构成新的商业权力。",
                f"{b('digital divide')}：技术收益与风险按阶层、年龄、地区、能力和教育不均分布。",
                f"{b('infrastructure')}：基础设施越顺畅，越会隐藏劳动、权力和维护成本。",
                f"{b('instrumental reason')}：工具理性关注效率和控制，可能遮蔽意义、情感、尊严和公共讨论。",
                f"{b('externality')}：技术成本被转嫁给环境、弱势群体、未来世代或不可见劳动者。",
            ],
        ),
        (
            "课程主题总表",
            [
                f"{b('生命')}：基因编辑、克隆、干细胞、医疗分流把生命变成可修复、可选择、可排序、可治理的对象。",
                f"{b('身体')}：钟表、工厂、铁路、学校、AI 工作流把身体纳入时间纪律、流程和绩效评价。",
                f"{b('数据')}：AI、NFC、社会信用把行为变成记录、预测和分类，便利与监控同时发生。",
                f"{b('环境')}：塑料和垃圾处理显示技术修复能力与消费制度之间的矛盾。",
                f"{b('亲密')}：AI 陪伴把情感支持平台化、数据化，改变责任、承诺与他者性的边界。",
                f"{b('公平')}：高价基因疗法、信用评分、数字支付、自动化创业都可能扩大能力差距。",
                f"{b('治理')}：许可、审查、透明、申诉、审计、比例原则、数据最小化、长期随访、公共参与。",
                f"{b('人文核心')}：不反对技术，而是追问技术服务何种生活、谁受益、谁承担代价、谁有权决定。",
            ],
        ),
    ],
]


PAGE1_FILLER = [
    [
        (
            "CRISPR 机制补充",
            [
                "细菌免疫三阶段：spacer acquisition 获取病毒片段；crRNA biogenesis 转录加工；interference 识别并切割入侵 DNA。",
                "递送方式：plasmid、mRNA、RNP、电转、病毒载体、脂质纳米颗粒；RNP 暴露时间短，潜在脱靶低。",
                "CCR5Δ32：自然存在的 32bp deletion 与 HIV 抵抗有关；但人工编辑 indel 不等于天然 Δ32。",
                "HBB/镰状细胞：突变导致异常血红蛋白；编辑策略包括直接修复 HBB 或重启 fetal hemoglobin。",
                "CAR-T 编辑：TRAC 位点定点插入可减少随机整合；PD-1 knockout 试图增强 T 细胞抗肿瘤活性。",
                "base editing：不造成 DSB 的单碱基改变；prime editing：pegRNA + nickase + reverse transcriptase，理论上更精确。",
                "伦理边界：严重单基因病治疗更容易被辩护；复杂性状增强最难正当化，因为基因-环境关系不清。",
            ],
        ),
        (
            "贺建奎后续与制度",
            [
                "CNN 材料：贺建奎出狱后提出阿尔茨海默相关基因编辑研究，引发科学家和生物伦理学者担忧。",
                "争议焦点从“能否编辑”扩展到“有过违规记录的科学家能否重新进入高风险领域”。",
                "中国监管强化：人类遗传资源管理、伦理审查办法、医学研究登记、重大伦理问题专家评议。",
                "全球共识倾向：体细胞治疗可在严格临床路径推进；生殖系编辑在安全性、必要性和社会共识不足前不应临床使用。",
            ],
        ),
    ],
    [
        (
            "克隆与干细胞补充",
            [
                "embryonic stem cells：多能性强，可分化为多类细胞；adult stem cells 限制较多但伦理争议较小。",
                "mitochondrial heteroplasmy：线粒体异质性是治疗性克隆/细胞替换需要处理的技术问题之一。",
                "immune compatibility：患者特异性 ntESC 理论上降低排斥，但制备过程复杂且成本高。",
                "IVF 剩余胚胎争议：捐赠研究、继续冷冻、销毁、植入之间涉及不同道德判断。",
                "患者倡议者立场：疾病痛苦本身具有道德重量；延迟研究也会造成可避免的伤害。",
                "宗教/保守生命伦理：早期胚胎具有从受精开始的生命连续性，因此不能被当作工具性材料。",
                "世俗生命伦理：胚胎道德地位可随发育阶段增强，许可研究不等于否认生命价值。",
            ],
        ),
        (
            "HFE 法律框架细节",
            [
                "licence：胚胎研究、储存、使用必须获得授权；authority 负责审查、记录、监督与处罚。",
                "consent：配子、胚胎使用和储存依赖明确同意；同意撤回会影响研究和生育安排。",
                "prohibited acts：无许可处理胚胎、超出授权目的、违反储存/使用规定都可能构成违法。",
                "法律逻辑：不是完全放任科研，也不是绝对禁止，而是用许可制把高风险研究纳入可追踪制度。",
            ],
        ),
    ],
    [
        (
            "AI/Agent 技术补充",
            [
                "tool use：模型调用搜索、数据库、代码、邮件、支付等外部工具后，输出不再只是文本，而会改变外部世界。",
                "memory：长期记忆提高连续性，也增加隐私、画像和操控风险。",
                "planning：任务拆解提高效率，但计划错误会在多步执行中累积。",
                "reflection：自我批评可减少错误，但也可能产生表面合理化，不能替代外部验证。",
                "prompt injection：外部网页/文档中的恶意指令诱导 agent 越权；需要权限隔离和内容可信度判断。",
                "least privilege：工具权限按任务最小化；高风险操作需要人工批准和日志。",
                "Anthropic 案例分析：有效利他主义、AI 安全、商业融资和快速扩张之间存在价值张力。",
            ],
        ),
        (
            "Constitutional AI 细节",
            [
                "CAI 第一阶段：模型根据原则批评自己的回答并修订；第二阶段：用 AI preference feedback 训练偏好模型。",
                "宪法原则来源可包括人权、伤害最小化、非歧视、隐私、儿童安全、避免违法等。",
                "优点：减少大量人工标注压力，让价值约束显性化；风险：原则选择本身带有文化和政治立场。",
                "harmlessness tax：过度安全可能降低有用性；CAI 目标是在 helpful 与 harmless 之间寻找更好平衡。",
                "red-team：主动寻找模型被诱导产生危险输出的路径；不是一次性测试，而是持续过程。",
            ],
        ),
    ],
    [
        (
            "NFC 场景补充",
            [
                "街边摊/超市：NFC 支付减少找零和排队，但也让消费决策更即时化。",
                "家庭场景：智能门锁、家电配网、蓝牙快速配对体现“触碰即连接”的低门槛。",
                "图书馆：RFID/NFC 可用于自助借还、盘点、定位和门禁防盗。",
                "医疗供应链：RFID 追踪药品库存、防篡改标签、减少人工贴标与盘点错误。",
                "NFC 与 RFID：NFC 是 RFID 的近距离、双向、交互化变体；RFID 距离更广，NFC 更强调安全交互。",
                "风险权衡：便利、安全、隐私、成本、用户习惯之间存在持续折中。",
            ],
        ),
        (
            "社会信用补充",
            [
                "Zhima Credit：平台化信用评分与消费、履约、身份、关系网络、行为偏好相关，体现商业评分逻辑。",
                "fair scoring：评分应避免以地区、年龄、职业、消费能力等变量间接歧视弱势群体。",
                "black box：个人若不知道分数如何形成，就难以纠错和改变处境。",
                "信用修复：允许纠错、申诉、履约后退出名单，是信用治理维持正当性的关键。",
                "联合惩戒的比例原则：失信行为与限制措施必须相称，不能无限扩展到无关生活领域。",
            ],
        ),
    ],
]


PAGE2_FILLER = [
    [
        (
            "塑料案例补充",
            [
                "深圳 110 米垃圾山：展示城市消费、土地压力和历史垃圾累积之间的关系。",
                "上海大型填埋场：塑料占比高，传统填埋无法消除长期环境负担。",
                "成都旧垃圾堆：复杂塑料与厨余混合，传统工艺难以处理，技术治理被迫面对历史遗留问题。",
                "二噁英：含氯塑料焚烧不当可能产生高毒性污染物，因此净化系统和温控条件至关重要。",
                "源头减量优先于末端治理；末端技术越强，越需要防止“继续消费无所谓”的意识形态。",
            ],
        ),
        (
            "时间与劳动补充",
            [
                "铁路时间：标准时间使长距离交通协调成为可能，也推动国家空间被统一计量。",
                "学校时间：课表、铃声、考试把学习划分为可管理单元。",
                "工厂时间：计时工资和流水线把劳动价值与时间单位绑定。",
                "数字时间：通知、日程、打卡、算法推荐让时间纪律从工厂扩展到日常生活。",
                "抵抗形式：缩短工时、休息权、慢生活、关闭通知、重新夺回注意力。",
            ],
        ),
    ],
    [
        (
            "生命政治补充",
            [
                "sovereign power：古典主权以生杀权显示力量；biopower 则以保护、管理、优化生命为名运作。",
                "security apparatus：不是禁止一切风险，而是计算、管理、分布和容忍一定风险。",
                "population as resource：人口与国家财富、劳动力、兵源、公共卫生和统计知识相连。",
                "normalization：通过平均值、标准、风险阈值定义正常/异常，并让个体主动调整自身。",
                "state of exception：紧急状态可能让非常措施常态化，健康、安全和秩序成为扩权理由。",
                "bare life：只剩生物性存活、缺乏政治资格和完整法律保护的生命状态。",
            ],
        ),
        (
            "阿甘本-福柯关系",
            [
                "论文主旨：阿甘本的生命政治不应只被看作反福柯，而是继承并扩展福柯的问题域。",
                "福柯侧重现代资本主义中的治理术、规训和人口管理；阿甘本把生命政治追溯到西方政治的主权结构。",
                "装置概念是桥梁：福柯用来分析权力-知识机制，阿甘本将其扩展为捕捉和塑造生命的一般机制。",
                "两者共同点：生命不再只是自然事实，而被政治、法律、医学、经济和技术持续征用。",
            ],
        ),
    ],
    [
        (
            "Hard Times 补充",
            [
                "Facts only：把儿童想象力、情感和经验压缩为可灌输事实，体现功利主义教育的贫乏。",
                "Coketown 的重复和烟尘表现工业现代性的单调、污染和去个性化。",
                "Stephen Blackpool：劳动者困境，显示工人被制度、婚姻法律、阶级和资本共同困住。",
                "Rachael：温柔、照护和道德稳定性，与工业社会的冷硬形成对照。",
                "小说不是反知识，而是反对把人缩减为事实、产量和计算结果。",
            ],
        ),
        (
            "铁路材料补充",
            [
                "司机的 Pointing and Calling：用手指和口呼把注意力外化，减少人为疏忽，是安全文化的一部分。",
                "调度员是铁路系统脑中心，需要实时处理客流、天气、线路状态和冲突。",
                "施工者面对高温、严寒、风沙、高原、冻土等环境，体现基础设施的身体代价。",
                "乘务员/检票员不仅检查车票，也维持秩序、回答问题、照顾乘客，体现技术系统中的情感劳动。",
                "不同国家铁路差异：高速自动化系统与苏丹等地区的慢速线路、人工判断、地方旅行文化形成对比。",
            ],
        ),
    ],
    [
        (
            "概念矩阵补充",
            [
                "autonomy：不是抽象选择，而要有信息、能力、替代方案和免于操控的条件。",
                "justice：包括分配正义、程序正义、承认正义；不仅问结果，也问谁被听见。",
                "dignity：反对把人只当作材料、数据、风险对象或效率单位。",
                "precautionary principle：面对不可逆和高不确定风险时，不能把证明伤害的责任完全推给受影响者。",
                "accountability：技术伤害发生后必须能追踪设计者、部署者、使用者和监管者责任。",
                "transparency：不是公开全部代码，而是让受影响者理解规则、后果和申诉路径。",
            ],
        ),
        (
            "课程关系网",
            [
                "基因编辑与克隆：生命可设计化；问题集中在胚胎、未来世代、治疗/增强、公平和不可逆性。",
                "AI 与社会信用：数据可预测化；问题集中在偏见、透明、操控、责任和治理边界。",
                "NFC 与钟表：便利和同步提高效率，也改变身体习惯、消费欲望和时间感。",
                "塑料与铁路：基础设施带来现代生活便利，同时隐藏环境成本和劳动成本。",
                "Hard Times 与生命政治：都警惕人被制度理性压缩为事实、身体、人口、分数或功能。",
                "人文维度：意义、情感、尊严、叙事、公共讨论和弱者位置，是技术评估不可删去的部分。",
            ],
        ),
    ],
]


PAGE1_DENSE = [
    [
        (
            "基因编辑知识串",
            [
                "CRISPR 时间线：1987 大肠杆菌重复序列被发现；2002 命名 CRISPR/Cas；2012 Doudna/Charpentier 证明可编程 Cas9；2016 人体临床试验加速。",
                "PAM 既是技术限制也是安全机制；不同 Cas 蛋白识别不同 PAM，影响可编辑位点范围。",
                "脱靶检测：GUIDE-seq、Digenome-seq、高通量测序；临床转化要求不仅看目标位点，也看全基因组结构变异。",
                "mosaicism 在胚胎编辑中特别危险，因为不同细胞携带不同编辑结果，出生个体无法通过一次检测完全确认风险。",
                "pleiotropy 指一个基因影响多个性状；CCR5 不只关系 HIV，也可能涉及免疫、炎症和神经系统功能。",
                "多基因性状如智力、身高、复杂疾病风险，不适合用单一基因编辑叙事解释。",
            ],
        ),
        (
            "基因伦理知识串",
            [
                "治疗性目标通常指减少严重疾病负担；增强性目标通常指提升正常范围内能力或选择偏好特征。",
                "未来儿童同意问题不是形式缺失，而是风险被不可逆地写入其身体和后代谱系。",
                "遗传病家庭的生育自主与未来儿童开放未来权之间存在冲突。",
                "高价疗法会把健康机会市场化；若公共医保覆盖，又涉及资源优先级和机会成本。",
                "全球治理难点：科学竞争、医疗旅游、国家监管差异、商业资本和患者希望交织。",
            ],
        ),
    ],
    [
        (
            "克隆知识串",
            [
                "SCNT 成功依赖核重编程；成年体细胞核必须被卵母细胞环境重新设定为早期胚胎状态。",
                "囊胚内细胞团可形成 ES cells；滋养层参与胎盘形成，是胚胎发育伦理讨论的重要分界。",
                "治疗性克隆的“治疗”并不意味着直接治愈，而是建立与患者匹配的细胞资源和疾病模型。",
                "卵母细胞需求使女性身体成为技术供应链的一部分，涉及风险、补偿、公平和知情同意。",
                "胚胎潜能论认为早期胚胎原则上可发展为人；关系论则强调发育环境、母体关系和社会承认。",
            ],
        ),
        (
            "医学伦理知识串",
            [
                "beneficence 关注减轻患者痛苦；non-maleficence 关注不制造不可接受风险；autonomy 关注同意；justice 关注资源和机会分配。",
                "医学研究伦理要求 scientific validity；没有可靠科学基础的高风险研究本身就不具备伦理正当性。",
                "患者希望不能自动压倒所有伦理限制，因为希望也可能被商业宣传和绝望处境塑造。",
                "临床前研究、动物实验、伦理审查、注册、长期随访共同构成高风险生物医学转化路径。",
                "禁止生殖性克隆并不必然禁止所有克隆技术；关键在目的、阶段、对象和制度控制。",
            ],
        ),
    ],
    [
        (
            "AI 知识串",
            [
                "预训练提供通用语言能力；微调塑造任务风格；对齐训练塑造偏好和边界；部署阶段还需要监控和反馈。",
                "模型“理解”与人类理解不同：它可操作语言模式，但不保证有经验、意向性或事实接触。",
                "数据污染会让评测失真；模型可能记住测试集、版权文本或个人信息。",
                "合成数据能扩充训练材料，也可能放大模型自身错误，形成 model collapse。",
                "多模态模型把文本、图像、音频、视频、动作接口连接起来，风险从语言输出扩展到感知和控制。",
                "AI 公司治理矛盾：安全承诺、商业增长、算力竞争、投资人压力和公众利益之间难以完全一致。",
            ],
        ),
        (
            "Agent 风险串",
            [
                "长链任务中的小错误会被后续步骤放大，尤其是订票、付款、发邮件、改代码、删文件等外部动作。",
                "工具调用日志是追责基础；没有日志，错误很难定位到模型、用户、工具或系统设计。",
                "人类批准节点适合放在不可逆、高成本、隐私敏感和法律后果明显的操作前。",
                "企业采用 agents 后，管理者可能把责任推给系统，把劳动风险推给员工或用户。",
                "AI 创业民主化不等于结果公平，因为数据、渠道、资本、算力和监管资源仍不均等。",
            ],
        ),
    ],
    [
        (
            "NFC/RFID 知识串",
            [
                "主动模式：发起设备产生 RF 场并供能；被动标签可无电池工作，适合交通卡、门禁卡、商品标签。",
                "卡模拟模式让手机变成卡，安全性取决于设备解锁、密钥存储、支付令牌和交易验证。",
                "relay attack 不一定需要破解加密，而是延长通信距离，让远处卡片像在现场一样响应。",
                "NFC 的短距离降低部分风险，但不能消除社工、恶意软件和用户操作诱导。",
                "无现金社会提高效率，也可能排斥老人、儿童、低收入者和没有智能手机的人。",
            ],
        ),
        (
            "数据治理知识串",
            [
                "数据最小化要求只收集完成目的所需数据；目的限定要求不得随意二次使用。",
                "去标识化不等于匿名化，多源数据可重新识别个人。",
                "算法歧视可来自训练数据、标签定义、特征选择、反馈循环和部署场景。",
                "自动化决策越影响教育、就业、金融、医疗、司法，就越需要解释、申诉和人工复核。",
                "信用系统把道德、法律和市场行为连接起来，容易把一次失误扩展为长期身份标签。",
            ],
        ),
    ],
]


PAGE2_DENSE = [
    [
        (
            "环境知识串",
            [
                "微塑料进入水体、土壤、食物链和人体组织，风险具有慢性、扩散和难回收特征。",
                "可降解塑料依赖特定温度、湿度和工业堆肥条件；在自然环境中不一定快速消失。",
                "焚烧发电把垃圾转化为能源，但也把污染治理压力转移到烟气净化、飞灰处理和邻避政治。",
                "EPR 让生产者承担产品废弃后的回收和处理责任，试图改变只追求销售的激励结构。",
                "绿色消费若只强调个人选择，会遮蔽企业包装、平台外卖、物流系统和政策设计的责任。",
            ],
        ),
        (
            "时间知识串",
            [
                "机械钟表把时间从地方性节律变成抽象统一尺度，是现代国家、市场和工业协作的基础。",
                "时间焦虑来自可量化目标不断增加：迟到、截止日期、效率排名、在线状态、即时回复。",
                "数字平台通过连续提醒和算法刷新占用碎片时间，使休息也被商业化。",
                "时间治理不是单纯压迫，也能带来公共协调；矛盾在于协调何时变成对生活的全面占有。",
                "现代自由常表现为“拥有自己的时间”，但工作、消费和技术平台不断争夺这一自由。",
            ],
        ),
    ],
    [
        (
            "生命政治知识串",
            [
                "规训权力制造 docile bodies：有用、服从、可训练、可比较的身体。",
                "生命权力通过统计知识发现人口规律，再通过卫生、保险、城市规划和风险管理干预。",
                "自由主义治理并非减少权力，而是通过选择、责任、自我管理让个体主动配合制度目标。",
                "疫情中的健康码、隔离、流调和风险区划展示生命政治与数字治理结合。",
                "基因编辑把生命在分子层面纳入治理；社会信用把行为在数据层面纳入治理。",
                "医学分流把生命价值转换为可比较指标，体现生命政治中的优化逻辑。",
            ],
        ),
        (
            "主权知识串",
            [
                "主权不是日常治理的反面；在危机、战争、疫情和安全事件中，主权决断会重新显现。",
                "例外状态的危险在于临时措施被常态化，个人权利长期处于悬置状态。",
                "集中营范式在阿甘本那里不是单一历史地点，而是法秩序内部排除生命的极端结构。",
                "赤裸生命概念揭示：仅被保护为生物性存活并不等于拥有完整政治主体地位。",
                "技术治理常以安全、健康、效率为理由扩权，因此需要持续检查例外是否被制度化。",
            ],
        ),
    ],
    [
        (
            "文学知识串",
            [
                "《艰难时世》把工业城市、学校、家庭和工厂连在一起，显示功利主义如何进入私人生活。",
                "Gradgrind 的教育排斥想象力，类似现代技术评估中只看指标而不看生活意义。",
                "Bounderby 的自我奋斗叙事遮蔽结构性不平等，是资本主义意识形态的一部分。",
                "Sissy 的价值不是反事实，而是补足事实无法包含的同情、经验和关系判断。",
                "Louisa 的痛苦说明没有情感教育的人难以处理欲望、责任和亲密关系。",
            ],
        ),
        (
            "基础设施知识串",
            [
                "基础设施具有背景性：它运行良好时被忽视，故障时才显现其社会重要性。",
                "铁路把空间压缩为时间表，改变城市、劳动力流动、区域经济和日常想象。",
                "高铁现代性依赖标准化轨道、信号系统、调度中心、维护制度和旅客纪律。",
                "基础设施劳动常被技术叙事遮蔽；工程师、施工者、维修工、司机、乘务员共同维持系统。",
                "地方铁路衰落不仅是交通问题，也关系社区记忆、产业转移和区域不平等。",
            ],
        ),
    ],
    [
        (
            "理论知识串",
            [
                "技术不是单个物件，而是 artifacts + institutions + users + regulations + imaginaries 的组合。",
                "社会建构论强调同一技术在不同群体中有不同意义，例如 NFC 是便利、风险、身份、消费和治理工具。",
                "风险社会中，风险常由技术系统制造，却由普通人、弱者或未来世代承担。",
                "监控资本主义把行为剩余转化为预测产品，再用推荐、广告和界面设计影响未来行为。",
                "数字鸿沟不仅是有没有设备，也包括技能、可负担性、语言、残障友好和制度可达性。",
            ],
        ),
        (
            "横向对照串",
            [
                "CRISPR 与 AI：都具有黑箱性和不确定性；前者改变生物代码，后者改变信息/行动流程。",
                "NFC 与社会信用：都依赖低摩擦数据流；一个优化支付/身份，一个优化治理/评分。",
                "塑料与 AI：都显示技术便利先行、外部性滞后显现；治理常追赶商业扩张。",
                "钟表与铁路：都把时间标准化，使现代协调成为可能，也把人纳入纪律。",
                "克隆与余数生命：都涉及生命价值能否被阶段、潜能、收益和资源指标衡量。",
                "文学与理论：Hard Times 的人被事实压扁，福柯/阿甘本的人被治理和主权捕获。",
            ],
        ),
    ],
]


def add_page_header(canvas, doc, title):
    canvas.saveState()
    width, height = A4
    canvas.setFont("STSong-Light", 7)
    canvas.setFillColor(colors.HexColor("#102131"))
    canvas.drawString(8 * mm, height - 7 * mm, title)
    canvas.setFont("STSong-Light", 5.5)
    canvas.drawRightString(width - 8 * mm, height - 7 * mm, f"A4 双面 | page {doc.page}")
    canvas.setStrokeColor(colors.HexColor("#78909c"))
    canvas.setLineWidth(0.25)
    canvas.line(8 * mm, height - 8.3 * mm, width - 8 * mm, height - 8.3 * mm)
    canvas.restoreState()


def build_pdf(extracted: dict[str, str]) -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    styles = make_styles()
    page_w, page_h = A4
    margin_x = 7.5 * mm
    top = 10.5 * mm
    bottom = 7 * mm
    gap = 2.2 * mm
    col_w = (page_w - margin_x * 2 - gap * 3) / 4
    frames = []
    for i in range(4):
        frames.append(
            Frame(
                margin_x + i * (col_w + gap),
                bottom,
                col_w,
                page_h - top - bottom,
                leftPadding=0,
                rightPadding=0,
                topPadding=0,
                bottomPadding=0,
            )
        )
    doc = BaseDocTemplate(
        str(OUT_PDF),
        pagesize=A4,
        leftMargin=margin_x,
        rightMargin=margin_x,
        topMargin=top,
        bottomMargin=bottom,
        title="科技与人文 双面A4超密cheatpaper",
        author="Codex",
    )
    doc.addPageTemplates(
        [
            PageTemplate(
                id="dense",
                frames=frames,
                onPage=lambda c, d: add_page_header(c, d, "科技与人文 知识点密集速查"),
            )
        ]
    )

    story = []

    story.append(para(styles["title"], "科技与人文 双面 A4 知识点密集速查"))
    for col_index, col_sections in enumerate(PAGE1_COLUMNS):
        for title, items in col_sections:
            story.extend(section(title, items, styles))
        for title, items in PAGE1_FILLER[col_index]:
            story.extend(section(title, items, styles))
        for title, items in PAGE1_DENSE[col_index]:
            story.extend(section(title, items, styles))
        if col_index < len(PAGE1_COLUMNS) - 1:
            story.append(FrameBreak())

    story.append(PageBreak())
    for col_index, col_sections in enumerate(PAGE2_COLUMNS):
        for title, items in col_sections:
            story.extend(section(title, items, styles))
        for title, items in PAGE2_FILLER[col_index]:
            story.extend(section(title, items, styles))
        for title, items in PAGE2_DENSE[col_index]:
            story.extend(section(title, items, styles))
        if col_index < len(PAGE2_COLUMNS) - 1:
            story.append(FrameBreak())

    doc.build(story)


def main() -> None:
    extracted = extract_all()
    build_pdf(extracted)
    print(OUT_PDF.resolve())
    print(NOTES_TXT.resolve())


if __name__ == "__main__":
    main()
